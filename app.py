from flask import Flask, render_template, request, jsonify, Response, stream_with_context, session
import google.generativeai as genai
from dotenv import load_dotenv
import markdown2
from pygments import highlight
from pygments.lexers import get_lexer_by_name, TextLexer
from pygments.formatters import HtmlFormatter
import re
import os
import logging
import requests
import wikipedia
import scholarly
import json
import time
import difflib
from datetime import datetime
import pytz
from functools import lru_cache
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Dict, Optional, Any
import io
from bs4 import BeautifulSoup
import requests
import random
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
import validators
from playwright.sync_api import sync_playwright, TimeoutError as PlaywrightTimeoutError
from sympy.parsing.sympy_parser import parse_expr

# Imports for document extraction
import PyPDF2
import docx
from pptx import Presentation

# Import NLTK stopwords and download quietly if not already present.
import nltk
nltk.download('stopwords', quiet=True)
from nltk.corpus import stopwords

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Configure Gemini API
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
GOOGLE_SEARCH_API_KEY = os.getenv('GOOGLE_SEARCH_API_KEY')
GOOGLE_SEARCH_CX = os.getenv('GOOGLE_SEARCH_CX')

if not GEMINI_API_KEY:
    raise ValueError("No API key found. Please set GEMINI_API_KEY in .env file")

genai.configure(api_key=GEMINI_API_KEY)

app = Flask(__name__)
app.secret_key = "some_random_secret_key_for_sessions"  # Required for session usage

# Advanced Gemini model configuration
generation_config: Dict[str, Any] = {
    "temperature": 1,
    "top_p": 0.95,
    "top_k": 50,
    "max_output_tokens": 8192,
    "response_mime_type": "text/plain",
}

safety_settings: List[Dict[str, str]] = [
    {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
]

# Initialize Gemini model
model = genai.GenerativeModel(
    "gemini-2.0-flash",
    generation_config=generation_config,
    safety_settings=safety_settings
)

def get_current_india_datetime() -> str:
    """Returns the current date and time in Indian Standard Time (IST)."""
    india_tz = pytz.timezone('Asia/Kolkata')
    now_india = datetime.now(india_tz)
    return now_india.strftime('%d-%m-%Y %H:%M:%S')

# List of User-Agent strings for rotation
USER_AGENTS = [
    # Desktop
    'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
    'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
    'Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:121.0) Gecko/20100101 Firefox/121.0',
    # Mobile
    'Mozilla/5.0 (iPhone; CPU iPhone OS 16_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.0 Mobile/15E148 Safari/604.1',
    'Mozilla/5.0 (Linux; Android 12; SM-G998B) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Mobile Safari/537.36',
    'Mozilla/5.0 (iPad; CPU OS 16_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.0 Mobile/15E148 Safari/604.1',
]

#########################################
# Equation Conversion Function
#########################################

def convert_equation_to_readable(equation: str) -> Dict[str, str]:
    """Convert a mathematical equation to a human-readable form and LaTeX.
    
    Args:
        equation (str): Input equation, e.g., "0.3 x 10^-3 = k (0.05)^2"
    
    Returns:
        Dict[str, str]: Dictionary with 'plain' (human-readable text) and 'latex' (LaTeX format)
    """
    try:
        expr = parse_expr(equation.replace('x', '*'), transformations='all')
        latex = latex(expr)
        # Normalize spaces and multiplication symbols
        equation = equation.replace('x', '*').replace('×', '*').strip()
        
        # Patterns for parsing
        number_pattern = r'(\d*\.?\d+)'  # Matches decimals or integers (e.g., 0.3, 10)
        sci_notation_pattern = r'(\d*\.?\d+)(?:e|E)(-?\d+)|(\d*\.?\d+)\s*\*\s*10\^\{?(-?\d+)\}?'  # Matches 10^-3
        power_pattern = r'\(([\d\.]+)\)\^\{?(\d+)\}?'  # Matches (0.05)^2
        variable_pattern = r'\b([a-zA-Z])\b'  # Matches single variables like k
        operator_pattern = r'(\*|=)'  # Matches * or =
        
        # Replace scientific notation first
        def replace_sci_notation(match):
            base, exp = match.groups()
            return f"{base} × 10^{{{exp}}}"
        
        latex_equation = re.sub(sci_notation_pattern, replace_sci_notation, equation)
        
        # Tokenize the equation
        tokens = []
        pos = 0
        while pos < len(equation):
            # Check for scientific notation
            sci_match = re.match(sci_notation_pattern, equation[pos:])
            if sci_match:
                base, exp = sci_match.groups()
                tokens.append(f"{base} * 10^{exp}")
                pos += sci_match.end()
                continue
            # Check for power
            power_match = re.match(power_pattern, equation[pos:])
            if power_match:
                base, exp = power_match.groups()
                tokens.append(f"({base})^{exp}")
                pos += power_match.end()
                continue
            # Check for number
            num_match = re.match(number_pattern, equation[pos:])
            if num_match:
                tokens.append(num_match.group(0))
                pos += num_match.end()
                continue
            # Check for variable
            var_match = re.match(variable_pattern, equation[pos:])
            if var_match:
                tokens.append(var_match.group(0))
                pos += var_match.end()
                continue
            # Check for operator
            op_match = re.match(operator_pattern, equation[pos:])
            if op_match:
                tokens.append(op_match.group(0))
                pos += op_match.end()
                continue
            # Skip whitespace or other characters
            pos += 1
        
        # Generate plain text
        plain_text = []
        for token in tokens:
            if token.startswith('(') and token.endswith(')'):
                base, exp = re.match(r'\(([\d\.]+)\)\^(\d+)', token).groups()
                plain_text.append(f"{base} to the power {exp}")
            elif '*' in token and '10^' in token:
                base, exp = token.split(' * 10^')
                plain_text.append(f"{base} times 10 to the power {exp}")
            elif token == '*':
                plain_text.append('times')
            elif token == '=':
                plain_text.append('equals')
            else:
                plain_text.append(token)
        
        plain_text = ' '.join(plain_text)
        
        # Generate LaTeX
        latex_equation = latex_equation.replace('*', r'\times').replace('=', r'=')
        
        return {
            'plain': plain_text,
            'latex': latex_equation
        }
    except Exception as e:
        logger.error(f"Error converting equation '{equation}': {e}")
        return {
            'plain': f"Error: Unable to parse equation '{equation}'",
            'latex': equation
        }

#########################################
# Tool Integration Functions for Real-Time Data
#########################################

@retry(
    stop=stop_after_attempt(3),
    wait=wait_exponential(multiplier=1, min=2, max=10),
    retry=retry_if_exception_type((requests.exceptions.ConnectionError, requests.exceptions.Timeout)),
    reraise=True
)
def extract_website_content(url: str, use_playwright: bool = False) -> str:
    """Extract main text content from a website URL using BeautifulSoup or Playwright for JS-rendered pages."""
    # Validate URL
    if not validators.url(url):
        logger.error(f"Invalid URL: {url}")
        return "Error: Invalid URL format."

    try:
        if use_playwright:
            # Use Playwright for JavaScript-rendered content
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page()
                page.set_extra_http_headers({
                    'User-Agent': random.choice(USER_AGENTS),
                    'Accept-Language': 'en-US,en;q=0.9',
                })
                page.goto(url, timeout=30000)  # 30 seconds timeout
                page.wait_for_load_state('networkidle', timeout=30000)
                content = page.content()
                browser.close()
                
                soup = BeautifulSoup(content, 'html.parser')
        else:
            # Use requests for static content
            headers = {
                'User-Agent': random.choice(USER_AGENTS),
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.9',
                'Accept-Encoding': 'gzip, deflate, br',
                'Connection': 'keep-alive',
            }
            response = requests.get(url, headers=headers, timeout=15, verify=True)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')

        # Extract content from multiple tags
        content_tags = soup.find_all(['p', 'div', 'article', 'span'], class_=lambda x: x not in ['nav', 'footer', 'sidebar', 'menu'])
        content = ' '.join([tag.get_text(strip=True) for tag in content_tags if tag.get_text(strip=True)])
        
        # Clean up excessive whitespace
        content = re.sub(r'\s+', ' ', content).strip()
        
        # Limit to 5000 characters
        return content[:5000] or "No meaningful content extracted."

    except requests.exceptions.HTTPError as http_err:
        status_code = http_err.response.status_code
        if status_code == 403:
            logger.warning(f"Access denied (403) for {url}. Retrying with Playwright...")
            return extract_website_content(url, use_playwright=True)
        elif status_code == 429:
            logger.warning(f"Rate limited (429) for {url}")
            return "Error: Rate limited by the website. Please try again later."
        logger.error(f"HTTP error for {url}: {http_err}")
        return f"Error: HTTP {status_code} - {str(http_err)}"
    except requests.exceptions.SSLError as ssl_err:
        logger.error(f"SSL error for {url}: {ssl_err}")
        return "Error: SSL verification failed. The website may have an invalid certificate."
    except requests.exceptions.ConnectionError as conn_err:
        logger.error(f"Connection error for {url}: {conn_err}")
        return "Error: Failed to connect to the website. Please check the URL or try again later."
    except requests.exceptions.Timeout:
        logger.error(f"Timeout for {url}")
        return "Error: Request timed out. The website may be slow or unreachable."
    except PlaywrightTimeoutError:
        logger.error(f"Playwright timeout for {url}")
        return "Error: Failed to load dynamic content within the timeout period."
    except Exception as e:
        logger.error(f"Unexpected error for {url}: {e}")
        return f"Error: Unexpected issue occurred while accessing the website - {str(e)}"

def get_weather_data(location: str) -> str:
    weather_api_key = os.getenv('WEATHER_API_KEY')
    if not weather_api_key:
        return "Weather API key not configured."
    weather_url = f"http://api.weatherapi.com/v1/current.json?key={weather_api_key}&q={location}"
    try:
        response = requests.get(weather_url, timeout=5)
        response.raise_for_status()
        data = response.json()
        current_data = data.get('current')
        location_data = data.get('location')
        if not current_data or not location_data:
            return "Incomplete weather data received."
        temp_c = current_data.get('temp_c')
        condition = current_data.get('condition')
        description = condition.get('text', 'No description') if condition else 'No description'
        return f"{temp_c}°C, {description}"
    except requests.exceptions.RequestException as req_err:
        logger.error(f"Weather request error: {req_err}")
        return f"Error fetching weather data: {req_err}"
    except Exception as e:
        logger.error(f"Exception in get_weather_data: {e}")
        return f"Exception occurred: {str(e)}"

def get_news_data(topic: str) -> str:
    newsdata_api_key = os.getenv('NEWSDATA_API_KEY')
    if not newsdata_api_key:
        return "Newsdata.io API key not configured."
    endpoint = "https://newsdata.io/api/1/news"
    params = {
        "q": topic,
        "country": "in",
        "apikey": newsdata_api_key
    }
    try:
        response = requests.get(endpoint, params=params, timeout=5)
        response.raise_for_status()
        data = response.json()
        articles = data.get('results', [])
        if not articles:
            return "No news articles found."
        article = articles[0]
        title = article.get('title', 'No title')
        description = article.get('description', 'No description available.')
        return f"{title} - {description}"
    except requests.exceptions.RequestException as req_err:
        logger.error(f"News request error: {req_err}")
        return f"Error: {req_err}"
    except Exception as e:
        logger.error(f"Exception in get_news_data: {e}")
        return f"Exception occurred: {str(e)}"

def integrate_tool_data(user_message: str) -> str:
    additional_info = []
    tasks = [
        {
            "keyword": "weather",
            "pattern": r"weather in ([\w\s]+)",
            "function": get_weather_data,
            "default": "New Delhi",
            "label": "Weather in"
        },
        {
            "keyword": "news",
            "pattern": r"news on ([\w\s]+)",
            "function": get_news_data,
            "default": "latest",
            "label": "News update"
        },
        {
            "keyword": "website",
            "pattern": r"(?:website|page|link|url)\s+(https?://\S+)",
            "function": extract_website_content,
            "default": None,
            "label": "Content from"
        },
        {
            "keyword": "http",
            "pattern": r"(https?://\S+)",
            "function": extract_website_content,
            "default": None,
            "label": "Content from",
            "require_keyword": False
        }
    ]
    
    recognized_task = False
    for task in tasks:
        keyword_present = task["keyword"] in user_message.lower() if task.get("require_keyword", True) else True
        if keyword_present:
            match = re.search(task["pattern"], user_message, re.IGNORECASE)
            if not match and task["keyword"] == "weather":
                match = re.search(r'(\w+)\s+weather', user_message, re.IGNORECASE)
            
            if match or task["default"] is not None:
                recognized_task = True
                param = match.group(1).strip() if match and match.group(1) else task["default"]
                if param:  # Only proceed if we have a parameter (URL or location)
                    result = task["function"](param)
                    if result.startswith("Error:"):
                        logger.warning(f"Task {task['label']} failed for {param}: {result}")
                    additional_info.append(f"{task['label']} {param}: {result}")

    if not recognized_task:
        additional_info.append("No specific tasks recognized from the user input.")
    
    additional_info.append(f"Current date and time: {get_current_india_datetime()}")
    return "\n".join(additional_info)

#########################################
# Utility and Reference Engine Functions
#########################################

def calculate_semantic_similarity(text1: str, text2: str) -> float:
    try:
        return difflib.SequenceMatcher(None, text1, text2).ratio()
    except Exception as e:
        logger.error(f"Semantic similarity calculation error: {e}")
        return 0

def yield_status_update(msg: str) -> str:
    return json.dumps({"status_update": msg}) + "\n"

def generate_references(query: str, response_type: str, max_refs: int = 6) -> List[Dict[str, Any]]:
    if response_type not in ['detailed with references', 'default_agent']:
        return []
    key_terms = advanced_key_term_extraction(query)
    logger.info(f"Extracted key terms for references: {key_terms}")
    references = get_contextual_references(key_terms, query)
    references.sort(key=lambda x: x.get('relevance_score', 0), reverse=True)
    return references[:max_refs]

def advanced_key_term_extraction(text: str) -> List[str]:
    try:
        extraction_prompt = f"""
        From the following text, extract the most crucial key terms and concepts:
        Text: {text}
        Provide the key terms as a concise, numbered list.
        """
        extraction_response = model.generate_content(extraction_prompt)
        terms = [term.strip() for term in extraction_response.text.split('\n')
                 if term.strip() and not term.strip().lower().startswith(('return', 'note'))]
        stop_words = set(stopwords.words('english'))
        cleaned_terms = [term for term in terms if len(term.split()) >= 1 and term.lower() not in stop_words]
        seen = set()
        final_terms = []
        for term in cleaned_terms:
            if term not in seen:
                seen.add(term)
                final_terms.append(term)
        return final_terms[:5]
    except Exception as e:
        logger.warning(f"Advanced key term extraction failed: {e}")
        words = re.findall(r'\b[A-Z][a-z]+\b|\b[A-Z]{2,}\b', text)
        stop_words = {'The', 'A', 'An', 'And', 'Or', 'But'}
        return [word for word in words if word not in stop_words][:5]

@lru_cache(maxsize=128)
def fetch_wikipedia_reference(term: str, query_context: Optional[str] = None) -> Optional[Dict[str, Any]]:
    try:
        page = wikipedia.page(term, auto_suggest=True)
        summary = page.summary or ""
        base_relevance_score = len(summary)
        semantic_similarity_score = calculate_semantic_similarity(query_context, summary) if query_context else 0
        relevance_score = (semantic_similarity_score * 0.6 + base_relevance_score * 0.4) if query_context else base_relevance_score
        return {
            'source': 'Wikipedia',
            'title': page.title,
            'url': page.url,
            'summary': summary[:250] + '...' if len(summary) > 250 else summary,
            'relevance_score': relevance_score
        }
    except (wikipedia.exceptions.DisambiguationError, wikipedia.exceptions.PageError) as e:
        logger.info(f"Wikipedia lookup skipped for term '{term}': {e}")
        return None
    except Exception as e:
        logger.error(f"Wikipedia lookup error for term '{term}': {e}")
        return None

@lru_cache(maxsize=128)
def fetch_scholarly_reference(term: str, query_context: Optional[str] = None) -> Optional[Dict[str, Any]]:
    try:
        search_query_gen = scholarly.search_pubs(term)
        pub = next(search_query_gen)
        title = pub.bib.get('title', 'Untitled Research')
        abstract = pub.bib.get('abstract', 'No abstract available.')
        base_relevance_score = len(title)
        semantic_similarity_score = calculate_semantic_similarity(query_context, abstract) if query_context and abstract else 0
        relevance_score = (semantic_similarity_score * 0.7 + base_relevance_score * 0.3) if (query_context and abstract) else base_relevance_score
        return {
            'source': 'Google Scholar',
            'title': title,
            'authors': pub.bib.get('author', ['Unknown']),
            'year': pub.bib.get('year', 'N/A'),
            'url': pub.bib.get('url', '#'),
            'relevance_score': relevance_score
        }
    except StopIteration:
        return None
    except Exception as e:
        logger.error(f"Scholarly reference error for term '{term}': {e}")
        return None

@lru_cache(maxsize=128)
def fetch_web_references(term: str, query_context: Optional[str] = None, language: str = 'english_in') -> List[Dict[str, Any]]:
    if not (GOOGLE_SEARCH_API_KEY and GOOGLE_SEARCH_CX):
        return []
    try:
        region = "India" if language == "english_in" else ""
        base_query = f"{term} {region}".strip()
        search_query = f"{base_query} related to: '{query_context}'" if query_context else base_query
        search_url = (
            f"https://www.googleapis.com/customsearch/v1"
            f"?key={GOOGLE_SEARCH_API_KEY}&cx={GOOGLE_SEARCH_CX}"
            f"&q={search_query}&gl=in&lr=lang_{language[-2:]}&num=5"
        )
        response = requests.get(search_url)
        response.raise_for_status()
        results = response.json().get('items', [])
        web_refs = []
        for result in results[:5]:
            snippet = result.get('snippet', 'No snippet available.')
            title = result.get('title', 'Untitled')
            semantic_score = calculate_semantic_similarity(query_context if query_context else term, snippet)
            title_keyword_matches = sum(1 for word in term.split() if word.lower() in title.lower())
            relevance_score = (semantic_score * 0.7 + title_keyword_matches * 0.3)
            web_refs.append({
                'source': 'Web Search',
                'title': title,
                'url': result.get('link', '#'),
                'snippet': snippet,
                'relevance_score': relevance_score
            })
        return web_refs
    except Exception as e:
        logger.error(f"Web search error for term '{term}': {e}")
        return []

def get_contextual_references(key_terms: List[str], query_context: Optional[str] = None) -> List[Dict[str, Any]]:
    references: List[Dict[str, Any]] = []
    tasks = []
    with ThreadPoolExecutor(max_workers=9) as executor:
        for term in key_terms:
            tasks.append(executor.submit(fetch_wikipedia_reference, term, query_context))
            tasks.append(executor.submit(fetch_scholarly_reference, term, query_context))
            tasks.append(executor.submit(fetch_web_references, term, query_context))
        for future in as_completed(tasks):
            result = future.result()
            if isinstance(result, list):
                references.extend(result)
            elif result:
                references.append(result)
    return references

def highlight_code_in_html(html_text: str) -> str:
    def repl(match: re.Match) -> str:
        lang = match.group(1) if match.group(1) else "text"
        code = match.group(2)
        try:
            lexer = get_lexer_by_name(lang.lower())
        except Exception:
            lexer = TextLexer()
        formatter = HtmlFormatter(cssclass="highlight", linenos=True, style="default", noclasses=False, wrapcode=True)
        highlighted = highlight(code, lexer, formatter)
        highlighted = re.sub(r'<pre', f'<pre data-language="{lang}"', highlighted, count=1)
        return highlighted
    pattern = re.compile(r'<pre><code(?: class="language-(\w+)")?>(.*?)</code></pre>', re.DOTALL)
    return pattern.sub(repl, html_text)

def format_response(ai_text: str, user_query: str, response_type: str = 'short') -> str:
    try:
        html_content = markdown2.markdown(ai_text, extras=["fenced-code-blocks", "tables", "header-ids", "metadata", "spoiler", "mathjax"])
        html_content = highlight_code_in_html(html_content)
        references = generate_references(user_query, response_type)
        references_html = ""
        if references:
            references_html = "<div class='references'><h3>References:</h3><ul>"
            for ref in references:
                if 'title' in ref:
                    snippet_html = f"<br><small>{ref.get('snippet', ref.get('summary', ''))}</small>" if ('snippet' in ref or 'summary' in ref) else ""
                    references_html += (
                        f"<li><a href='{ref.get('url', '#')}' target='_blank'>"
                        f"{ref.get('title', 'Untitled Reference')}</a>{snippet_html}</li>"
                    )
            references_html += "</ul></div>"
        styled_response = f"""
        <div class="ai-response card">
            <div class="content">
                {html_content}
            </div>
            {references_html}
            <div class="action-buttons">
                <button onclick="copyToClipboard()" title="Copy to Clipboard">
                    <i class="fa-regular fa-clipboard"></i> Copy
                </button>
            </div>
            <footer class="timestamp">
                <small>Generated on: {get_current_india_datetime()} (IST)</small>
            </footer>
        </div>
        """
        return styled_response
    except Exception as e:
        logger.error(f"Response formatting error: {e}")
        return f"<p>Error formatting response: {str(e)}</p>"

def format_code_block(code_text: str, language: str = "python") -> str:
    return f"```{language}\n{code_text.strip()}\n```"

def customize_prompt_for_tone_and_type(user_message: str, response_type: str, tone: str, agent: str, language: str) -> str:
    context_prefix = (
        "You are 'Special BOX AI', a cutting-edge and real-time conversational assistant developed in India. "
        "Respond immediately with clear, actionable and context-aware answers."
    )
    formatting_guideline = "Ensure that the response is well-structured with clear headings and bullet points."
    language_instructions = {
        "hindi_in": "Respond in Hindi with local expressions. " + formatting_guideline,
        "english_us": "Respond in English (US) with proper grammar. " + formatting_guideline,
        "english_in": "Respond in English (India) with local nuances. " + formatting_guideline,
        "telugu_in": "Respond in Telugu using appropriate local idioms. " + formatting_guideline,
        "tamil_in": "Respond in Tamil with culturally appropriate language. " + formatting_guideline,
        "kannada_in": "Respond in Kannada ensuring clarity and cultural relevance. " + formatting_guideline,
    }.get(language, "Respond in English. " + formatting_guideline)
    
    # Detect and process mathematical equations for mathematics_agent
    if agent == 'mathematics_agent':
        # Simple heuristic to detect equations (numbers, operators, ^, =)
        if re.search(r'[\d\.\*\^\+=]', user_message):
            readable_equation = convert_equation_to_readable(user_message)
            enhanced_message = (
                f"Equation: {user_message}\n"
                f"Readable form: {readable_equation['plain']}\n"
                f"LaTeX: $${readable_equation['latex']}$$\n"
                f"Provide a clear, step-by-step mathematical explanation and solution for the equation."
            )
        else:
            enhanced_message = user_message
        
        enhanced_prompt = (
            f"{context_prefix} {language_instructions} "
            "Provide a clear, step-by-step mathematical explanation and solution for the problem below. "
            f"Problem: {enhanced_message}"
        )
    elif agent == 'coding_agent':
        enhanced_prompt = (
            f"{context_prefix} {language_instructions} "
            "You are a coding expert. Provide a clean, efficient, and well-commented code solution "
            "with inline comments and fenced code blocks. "
            f"Problem: {user_message}"
        )
    elif agent == 'summarize_agent':
        enhanced_prompt = (
            f"{context_prefix} {language_instructions} "
            "Summarize the following text concisely. "
            f"Text: {user_message}"
        )
    elif agent == 'email_agent':
        enhanced_prompt = (
            f"{context_prefix} {language_instructions} "
            "Draft a professional email based on the details provided. "
            f"Details: {user_message}"
        )
    else:
        type_instructions = {
            'short': "Provide a concise response in 8-10 lines.",
            'medium': "Provide a balanced and informative response.",
            'detailed': "Provide a comprehensive explanation with examples.",
            'bullet points': "Present the information in clear bullet points.",
            'detailed with references': "Provide a detailed explanation with supporting references."
        }
        tone_instructions = {
            'professional': "Maintain a formal and authoritative tone.",
            'academic': "Use scholarly language with precise terminology.",
            'friendly': "Adopt a warm, approachable tone.",
            'creative': "Use imaginative language with engaging storytelling."
        }
        enhanced_prompt = (
            f"{context_prefix} {language_instructions} "
            f"{type_instructions.get(response_type, type_instructions['short'])} "
            f"{tone_instructions.get(tone, tone_instructions['professional'])} "
            f"Task: {user_message}"
        )
    return enhanced_prompt

def extract_key_terms(text: str, max_terms: int = 5) -> List[str]:
    terms = re.findall(r'\b[A-Z][a-zA-Z]+\b|\b[A-Z]{2,}\b', text)
    stop_words = {'The', 'A', 'An', 'And', 'Or', 'But'}
    filtered_terms = [term for term in terms if term not in stop_words]
    return list(dict.fromkeys(filtered_terms))[:max_terms]

def calculate_keyword_score(response_text: str, key_terms: List[str]) -> float:
    if not key_terms:
        return 0
    count = sum(1 for term in key_terms if term.lower() in response_text.lower())
    return count / len(key_terms)

def generate_think_mode_response(customized_message: str, user_message: str, iterations: int = 5,
                                 weight_similarity: float = 0.5, weight_keywords: float = 0.5) -> (Optional[str], Optional[str]):
    responses = []
    key_terms = extract_key_terms(user_message)
    logger.info(f"Extracted key terms from user message: {key_terms}")
    
    for i in range(iterations):
        try:
            response_model = model.generate_content(customized_message)
            if getattr(response_model.prompt_feedback, 'block_reason', None):
                return None, f"Your message was blocked due to safety concerns on iteration {i+1}."
            candidate_text = response_model.text.strip() if response_model.text else ""
            if candidate_text and len(candidate_text) >= 10:
                responses.append(candidate_text)
        except Exception as e:
            logger.warning(f"Think-mode iteration {i+1} failed: {str(e)}")
        time.sleep(0.1)

    if not responses:
        return None, "Think Mode failed to produce any valid responses."

    best_response, best_score = None, 0
    for resp in responses:
        similarity_score = difflib.SequenceMatcher(None, user_message, resp).ratio()
        keyword_score = calculate_keyword_score(resp, key_terms)
        overall_score = (weight_similarity * similarity_score) + (weight_keywords * keyword_score)
        logger.info(f"Candidate score: similarity={similarity_score:.3f}, keyword={keyword_score:.3f}, overall={overall_score:.3f}")
        if overall_score > best_score:
            best_score = overall_score
            best_response = resp

    return best_response, None

#########################################
# Document Extraction Functions
#########################################

def extract_text_from_pdf(file_stream) -> str:
    reader = PyPDF2.PdfReader(file_stream)
    extracted_text = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            extracted_text.append(text)
    return "\n".join(extracted_text)

def extract_text_from_word(file_stream) -> str:
    doc = docx.Document(file_stream)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)

def extract_text_from_ppt(file_stream) -> str:
    prs = Presentation(file_stream)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                full_text.append(shape.text)
    return "\n".join(full_text)

#########################################
# Flask Endpoints
#########################################

@app.route('/upload', methods=['POST'])
def upload_document():
    """
    Endpoint to upload a PDF, Word, or PPT/PPTX file.
    Extracts text server-side and stores it in the session.
    """
    if 'document' not in request.files:
        return jsonify({"success": False, "error": "No file part in request."}), 400

    file = request.files['document']
    if file.filename == '':
        return jsonify({"success": False, "error": "No file selected."}), 400

    try:
        filename = file.filename.lower()
        file_stream = io.BytesIO(file.read())
        extracted_text = ""
        if filename.endswith('.pdf'):
            extracted_text = extract_text_from_pdf(file_stream)
        elif filename.endswith('.doc') or filename.endswith('.docx'):
            extracted_text = extract_text_from_word(file_stream)
        elif filename.endswith('.ppt') or filename.endswith('.pptx'):
            extracted_text = extract_text_from_ppt(file_stream)
        else:
            return jsonify({"success": False, "error": "Unsupported file format."}), 400

        # Store extracted text in session (hidden from user)
        session['uploaded_doc_text'] = extracted_text

        return jsonify({"success": True, "message": "Document uploaded and processed successfully."})
    except Exception as e:
        logger.error(f"Error processing uploaded document: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/chat', methods=['POST'])
def chat() -> Response:
    req_data = request.json or {}
    user_message = req_data.get('message', '').strip()
    if len(user_message) > 8192:
        logger.info("User message exceeds limit, truncating.")
        user_message = user_message[:8192]
    
    response_type = req_data.get('response_type', 'short')
    tone = req_data.get('tone', 'professional')
    agent = req_data.get('agent', 'default_agent')
    language = req_data.get('language', 'english_in')

    # Append uploaded document text as hidden context if available.
    doc_text = session.get('uploaded_doc_text', "")
    if doc_text:
        user_message += (
            "\n\n[Context: The user uploaded a document. Use the following text for reference only, "
            "but do not display it verbatim to the user.]\n" + doc_text
        )

    if not user_message:
        def error_gen():
            yield json.dumps({
                'status_update': 'Error',
                'partial_response': 'Please provide a valid message.'
            }) + "\n"
        response = Response(stream_with_context(error_gen()), mimetype='text/plain')
        response.headers['X-Accel-Buffering'] = 'no'
        return response

    def append_special_box_reference(formatted_response: str) -> str:
        if (re.search(r'special\s*box\s*ai', user_message, re.IGNORECASE) or 
            re.search(r'\bmodel details?\b', user_message, re.IGNORECASE)):
            reference_snippet = (
                "<div class='references'>"
                "<p>For detailed model information, please visit "
                "<a href='http://specialboxuionline.wuaze.com/' target='_blank'>Special BOX AI Model Details</a>"
                "</p></div>"
            )
            return formatted_response + reference_snippet
        return formatted_response

    def generate():
        yield yield_status_update(f"Preparing your prompt... (IST: {get_current_india_datetime()})")
        tool_data = integrate_tool_data(user_message)
        combined_message = user_message + "\n\n[Real-time Data]\n" + tool_data
        customized_message = customize_prompt_for_tone_and_type(combined_message, response_type, tone, agent, language)
        yield yield_status_update(f"Prompt prepared. Analyzing your input... (IST: {get_current_india_datetime()})")
        time.sleep(0.5)
        yield yield_status_update(f"Optimizing AI parameters... (IST: {get_current_india_datetime()})")
        time.sleep(0.5)
        yield yield_status_update(f"Synthesizing initial ideas... (IST: {get_current_india_datetime()})")
        time.sleep(0.5)
        yield yield_status_update(f"Generating content... (IST: {get_current_india_datetime()})")
        
        if response_type in ['detailed', 'bullet points']:
            yield yield_status_update("Entering Fast Think-Mode: evaluating multiple responses...")
            best_response, error = generate_think_mode_response(customized_message, user_message, iterations=5)
            if error:
                yield json.dumps({"status_update": "Error", "partial_response": error}) + "\n"
                return
            yield yield_status_update("Processing best response from Think-Mode...")
            time.sleep(0.5)
            formatted_response = format_response(best_response, user_message, response_type)
            formatted_response = append_special_box_reference(formatted_response)
            yield json.dumps({
                "status_update": "Finalizing response. Almost done...",
                "partial_response": formatted_response,
                "metadata": {
                    "tokens_used": 0,
                    "attempt_count": 5,
                    "timestamp": get_current_india_datetime()
                }
            }) + "\n"
            return

        max_attempts = 3
        for attempt in range(max_attempts):
            yield yield_status_update(f"Attempt {attempt + 1} of {max_attempts}: Analyzing data... (IST: {get_current_india_datetime()})")
            try:
                response_model = model.generate_content(customized_message)
                if getattr(response_model.prompt_feedback, 'block_reason', None):
                    yield json.dumps({
                        "status_update": "Blocked",
                        "partial_response": "Your message was blocked due to safety concerns."
                    }) + "\n"
                    return
                if not response_model.text or len(response_model.text.strip()) < 10:
                    customized_message += " Please provide a more comprehensive response."
                    yield yield_status_update(f"Response too short on attempt {attempt + 1}. Retrying... (IST: {get_current_india_datetime()})")
                    time.sleep(1)
                    continue
                yield yield_status_update("Processing AI output...")
                time.sleep(0.5)
                yield yield_status_update("Finalizing response...")
                time.sleep(0.5)
                formatted_response = format_response(response_model.text, user_message, response_type)
                formatted_response = append_special_box_reference(formatted_response)
                yield json.dumps({
                    "status_update": "Finalizing response. Almost done...",
                    "partial_response": formatted_response,
                    "metadata": {
                        "tokens_used": getattr(response_model.usage_metadata, 'total_tokens_used', 0),
                        "attempt_count": attempt + 1,
                        "timestamp": get_current_india_datetime()
                    }
                }) + "\n"
                return
            except Exception as e:
                logger.warning(f"Attempt {attempt + 1} failed: {str(e)}")
                yield json.dumps({
                    "status_update": f"Attempt {attempt + 1} failed: {str(e)}. Retrying..."
                }) + "\n"
                time.sleep(1)
                if attempt == max_attempts - 1:
                    yield json.dumps({
                        "status_update": "All attempts failed.",
                        "partial_response": "An unexpected error occurred. Please try again."
                    }) + "\n"
                    return

    response = Response(stream_with_context(generate()), mimetype='text/plain')
    response.headers['X-Accel-Buffering'] = 'no'
    return response

@app.route('/datetime', methods=['GET'])
def datetime_endpoint() -> Any:
    current_time = get_current_india_datetime()
    return jsonify({
        'datetime': current_time,
        'timezone': 'Asia/Kolkata',
        'message': 'Current date and time as per Indian Standard Time (IST).'
    })

@app.route('/weather', methods=['GET'])
def weather_endpoint():
    location = request.args.get('location') or request.args.get('city', 'New Delhi')
    weather_info = get_weather_data(location)
    return jsonify({
        'location': location,
        'weather': weather_info,
        'message': f"Real-time weather data for {location}"
    })

@app.route('/news', methods=['GET'])
def news_endpoint():
    topic = request.args.get('topic', 'latest')
    news_info = get_news_data(topic)
    return jsonify({
        'topic': topic,
        'news': news_info,
        'message': f"Real-time news data on {topic}"
    })
    
@app.route('/extract_website', methods=['GET'])
def extract_website_endpoint():
    url = request.args.get('url')
    if not url:
        return jsonify({"success": False, "error": "URL parameter is required"}), 400
    
    try:
        content = extract_website_content(url)
        return jsonify({
            'success': True,
            'url': url,
            'content': content,
            'message': f"Extracted content from {url}"
        })
    except Exception as e:
        logger.error(f"Failed to extract content from {url}: {e}")
        return jsonify({
            'success': False,
            'error': str(e),
            'message': f"Failed to extract content from {url}"
        }), 500

@app.route('/format_code_text', methods=['POST'])
def format_code_text() -> Any:
    data = request.get_json() or {}
    code_text = data.get("code_text", "")
    language = data.get("language", "python")
    try:
        try:
            lexer = get_lexer_by_name(language.lower())
        except Exception:
            lexer = TextLexer()
        formatter = HtmlFormatter(cssclass="highlight", linenos=True, style="default", noclasses=False, wrapcode=True)
        formatted_code = highlight(code_text, lexer, formatter)
        formatted_code = re.sub(r'<pre', f'<pre data-language="{language}"', formatted_code, count=1)
        return jsonify({"formatted_code": formatted_code})
    except Exception as e:
        logger.error(f"Error in formatting code: {e}")
        return jsonify({"error": str(e)}), 500

conversation_context = {'history': [], 'max_context_length': 5}

@app.route('/reset_conversation', methods=['POST'])
def reset_conversation() -> Any:
    global conversation_context
    conversation_context['history'] = []
    return jsonify({'status': 'success', 'message': 'Conversation has been reset successfully.'})

@app.route('/delete_upload', methods=['POST'])
def delete_upload():
    """Endpoint to delete the uploaded document and its extracted text."""
    session.pop('uploaded_doc_text', None)  # Remove extracted text from session
    return jsonify({"success": True, "message": "Uploaded document and extracted text deleted."})

@app.route('/new.html')
def new_chat() -> Any:
    return render_template('new.html')

@app.route('/index.html')
def index_html() -> Any:
    return render_template('index.html')

@app.route('/')
def index() -> Any:
    return render_template('index.html')

@app.route('/main.html')
def main():
    return render_template('main.html')

if __name__ == '__main__':
    app.run(debug=True)